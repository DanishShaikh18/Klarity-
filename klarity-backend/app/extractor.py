#extractor.py

import os
import fitz
import pdfplumber
from pdf2image import convert_from_path
import pytesseract
from docx import Document as DocxDocument
from pptx import Presentation
from langchain_community.document_loaders import PyMuPDFLoader
from langchain_core.documents import Document


# ---------------------------------------------------------
# 🔍 PDF HELPERS
# ---------------------------------------------------------

def is_page_scanned(page):
    """Check if a PDF page likely needs OCR."""
    text = page.get_text().strip()
    return len(text) < 25


def ocr_pdf_pages(pdf_path, pages_to_ocr):
    """Apply OCR only to scanned pages (efficient)."""
    images = convert_from_path(pdf_path, dpi=250)
    ocr_results = {}

    for idx in pages_to_ocr:
        text = pytesseract.image_to_string(images[idx], lang='eng').strip()
        if text:
            ocr_results[idx + 1] = text

    return ocr_results


def extract_pdf_tables(pdf_path):
    """Extract tables from PDF as separate documents."""
    tables = []

    with pdfplumber.open(pdf_path) as pdf:
        for page_num, page in enumerate(pdf.pages, start=1):
            for table in page.extract_tables():
                tables.append(Document(
                    page_content=str(table),
                    metadata={"type": "table", "page": page_num, "source": pdf_path}
                ))

    return tables


# ---------------------------------------------------------
# 📄 PDF MAIN EXTRACTOR
# ---------------------------------------------------------

def extract_pdf(pdf_path):
    """Primary: LangChain PDF extract → OCR fallback → tables."""
    loader = PyMuPDFLoader(pdf_path)
    docs = loader.load()  # one doc per page

    pdf = fitz.open(pdf_path)
    scanned_pages = [i for i, page in enumerate(pdf) if is_page_scanned(page)]

    if scanned_pages:
        ocr_text = ocr_pdf_pages(pdf_path, scanned_pages)
        for d in docs:
            page_no = d.metadata.get("page")
            if page_no in ocr_text:
                d.page_content = ocr_text[page_no]

    tables = extract_pdf_tables(pdf_path)

    return docs + tables


# ---------------------------------------------------------
# 📘 WORD EXTRACTOR (.docx)
# ---------------------------------------------------------

def extract_docx(docx_path):
    doc = DocxDocument(docx_path)
    content = []

    for para in doc.paragraphs:
        if para.text.strip():
            content.append(para.text.strip())

    # include tables
    for table in doc.tables:
        for row in table.rows:
            row_data = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if row_data:
                content.append(" | ".join(row_data))

    return "\n".join(content)


# ---------------------------------------------------------
# 🖼️ PPT EXTRACTOR (.pptx)
# ---------------------------------------------------------

def extract_ppt(pptx_path):
    prs = Presentation(pptx_path)
    output = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        output.append(f"--- Slide {slide_num} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                output.append(shape.text.strip())

        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                output.append(f"Notes: {notes}")

    return "\n".join(output)


# ---------------------------------------------------------
# 🎯 MASTER ROUTER
# ---------------------------------------------------------

def extract_file_content(file_path):
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        return extract_pdf(file_path)

    elif ext == ".docx":
        return extract_docx(file_path)

    elif ext == ".pptx":
        return extract_ppt(file_path)

    else:
        raise ValueError(f"❌ Unsupported file format: {ext}")


# ---------------------------------------------------------
# 🧪 TEST (Remove Later)
# ---------------------------------------------------------







